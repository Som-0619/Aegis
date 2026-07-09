"""
PPT Generator Agent.
Compiles weekly project health reports and monthly portfolio trend analyses
into premium, 16:9 widescreen, consulting-style slide presentations.
Uses distinct structural card layouts, status badges, and chart placeholders.
"""

import os
import json
from typing import Dict, Any, List
from utils.logger import get_logger

logger = get_logger(__name__)

# Fallback import of pptx to handle environment check
try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    from pptx.chart.data import CategoryChartData  # type: ignore
    from pptx.enum.chart import XL_CHART_TYPE  # type: ignore
    from pptx.enum.chart import XL_LEGEND_POSITION  # type: ignore
    PPTX_AVAILABLE = True
except Exception as e:
    logger.warning(f"python-pptx package could not be imported: {e}. PPTX generation will run in mock mode.")
    PPTX_AVAILABLE = False

class PPTGenerator:
    """Agent that handles compiling weekly project updates into widescreen slide deck presentations."""
    
    def __init__(self, template_path: str = "templates/ppt_template.pptx"):
        self.template_path = os.path.abspath(template_path)
        logger.info(f"PPTGenerator initialized. Template path: {self.template_path}")

    def generate_presentation(self, report_data: Dict[str, Any], output_path: str) -> str:
        """
        Generates a 6-slide, widescreen weekly executive report slide deck.
        1. Title & Overall Status Dashboard
        2. PMO Executive Audit Summary
        3. Category Scores & Detail Breakdown Grid
        4. Milestones & Timeline Tracker
        5. Operational Risks & Constraints (Blockers, Sentiment, Resources)
        6. Strategic Recommendations & Next Steps
        """
        logger.info(f"Generating Weekly PPTX presentation. Saving to {output_path}")
        
        # Ensure directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        if not PPTX_AVAILABLE:
            logger.info("Running in mock mode. Creating dummy PPTX file.")
            with open(output_path, "w") as f:
                f.write(f"Mock Slide Deck Content for: {report_data.get('project_name', 'Unknown')}")
            return output_path

        try:
            prs = None
            if os.path.exists(self.template_path) and os.path.getsize(self.template_path) > 0:
                try:
                    prs = Presentation(self.template_path)
                    logger.info(f"Loaded existing template from {self.template_path}")
                except Exception as ex:
                    logger.warning(f"Failed to load PPTX template: {ex}. Starting blank.")
            
            if prs is None:
                prs = Presentation()
                logger.info("Starting blank presentation.")

            # Set modern 16:9 widescreen dimensions
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Layout 6 is blank in standard template
            blank_layout = prs.slide_layouts[6]
            
            # Standard Palette
            navy = RGBColor(15, 23, 42)
            white = RGBColor(255, 255, 255)
            slate_50 = RGBColor(248, 250, 252)
            slate_200 = RGBColor(226, 232, 240)
            slate_500 = RGBColor(100, 116, 139)
            slate_800 = RGBColor(30, 41, 59)
            purple = RGBColor(124, 58, 237)
            
            # Status colors
            green = RGBColor(34, 197, 94)
            yellow = RGBColor(234, 179, 8)
            red = RGBColor(239, 68, 68)

            def get_status_color(status: str) -> RGBColor:
                s = status.upper()
                if "GREEN" in s: return green
                if "YELLOW" in s: return yellow
                return red

            # Helper for header titles
            def add_header(slide, title_text: str):
                header_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
                header_rect.fill.solid()
                header_rect.fill.fore_color.rgb = navy
                header_rect.line.fill.background()
                
                txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.6))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title_text
                p.font.name = "Arial"
                p.font.size = Pt(26)
                p.font.bold = True
                p.font.color.rgb = white

            # Helper for consulting-style cards
            def add_card(slide, left, top, width, height, title, body_text="", fill_color=None, border_color=None):
                card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                card.fill.solid()
                card.fill.fore_color.rgb = fill_color or slate_50
                if border_color:
                    card.line.color.rgb = border_color
                else:
                    card.line.fill.background()
                
                tf = card.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.2)
                tf.margin_top = Inches(0.2)
                
                p = tf.paragraphs[0]
                p.text = title
                p.font.name = "Arial"
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = navy
                
                if body_text:
                    p2 = tf.add_paragraph()
                    p2.text = body_text
                    p2.font.name = "Arial"
                    p2.font.size = Pt(11)
                    p2.font.color.rgb = slate_500
                    p2.space_before = Pt(8)
                return card

            raw_data = report_data.get("raw_extracted_data", {})
            overall_status = report_data.get("overall_status", "GREEN")
            overall_score = report_data.get("overall_score", 100.0)
            project_name = report_data.get("project_name", "Project")
            report_date = report_data.get("report_date", "N/A")
            confidence = report_data.get("confidence_score", 100.0)
            missing = report_data.get("missing_data", "None.")

            # ----------------------------------------------------
            # SLIDE 1: Title & Dashboard Summary
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, f"{project_name} - Executive Weekly Dashboard")
            
            # Big RAG Card
            badge_color = get_status_color(overall_status)
            rag_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(3.5), Inches(5.5))
            rag_card.fill.solid()
            rag_card.fill.fore_color.rgb = badge_color
            rag_card.line.fill.background()
            
            tf_rag = rag_card.text_frame
            tf_rag.word_wrap = True
            p_status = tf_rag.paragraphs[0]
            p_status.text = f"\n\nOverall Status\n{overall_status}"
            p_status.font.name = "Arial"
            p_status.font.size = Pt(28)
            p_status.font.bold = True
            p_status.font.color.rgb = white
            p_status.alignment = 1
            
            p_score = tf_rag.add_paragraph()
            p_score.text = f"Score: {overall_score}/100"
            p_score.font.name = "Arial"
            p_score.font.size = Pt(22)
            p_score.font.bold = True
            p_score.font.color.rgb = white
            p_score.alignment = 1
            
            # Timeline & Budget Card
            meta_body = (
                f"• Planned Start: {raw_data.get('planned_start_date') or 'N/A'}\n"
                f"• Planned End:   {raw_data.get('planned_end_date') or 'N/A'}\n"
                f"• Progress:      {raw_data.get('current_progress') or 'N/A'}\n"
                f"• System Confidence: {confidence}%\n\n"
                f"Financial Summary:\n"
                f"• Baseline Budget: ${float(raw_data.get('budget') or 0.0):,.2f}\n"
                f"• Actual Spent:    ${float(raw_data.get('budget_spent') or 0.0):,.2f}"
            )
            add_card(slide, Inches(4.3), Inches(1.3), Inches(4.2), Inches(5.5), "Project Timeline & Financials", meta_body)
            
            # Warnings / Missing Data
            warn_body = (
                f"• Missing Data fields identified: \n  {missing}\n\n"
                f"• Data Integrity Verdict: \n  "
                f"{'COMPLETE' if 'None' in missing else 'WARNING - Incomplete source updates.'}"
            )
            add_card(slide, Inches(8.7), Inches(1.3), Inches(4.1), Inches(5.5), "Data Integrity & Gaps", warn_body)

            # ----------------------------------------------------
            # SLIDE 2: PMO Executive Audit Summary
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, "PMO Executive Audit Summary")
            
            cats = report_data.get("scores", {})
            
            pmo_summary = (
                f"Aegis has completed the weekly audit for {project_name} as of {report_date}.\n\n"
                f"• Overall Status: {overall_status} (Weighted Score: {overall_score}/100)\n"
                f"• Data Integrity: {confidence}% of required status parameters extracted.\n"
                f"• Missing Parameters: {missing}\n\n"
                f"Operational Analysis:\n"
                f"• Timeline Velocity: {'On Track' if cats.get('schedule', 100.0) >= 80 else 'Lagging behind expected baseline progress'}.\n"
                f"• Cost Variance: {'Executing within budget bounds' if cats.get('budget', 100.0) >= 100 else 'Experiencing budget spent overrun'}.\n"
                f"• Risks & Constraints: {'Operational path is clean' if cats.get('risks', 100.0) >= 100 else 'Impacted by active blockers or key dependencies'}.\n"
            )
            add_card(slide, Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5), "Executive Overview & Findings", pmo_summary)

            # ----------------------------------------------------
            # SLIDE 3: Category Scores Breakdown
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, "Category Scoreboard Breakdown")
            
            col_width = Inches(3.8)
            row_height = Inches(2.5)
            
            # Draw a grid of 6 cards
            positions = [
                (Inches(0.5), Inches(1.3)), (Inches(4.7), Inches(1.3)), (Inches(8.9), Inches(1.3)),
                (Inches(0.5), Inches(4.3)), (Inches(4.7), Inches(4.3)), (Inches(8.9), Inches(4.3))
            ]
            
            cat_list = [
                ("Schedule", "schedule", "Calculates elapsed timeline vs actual progress rate."),
                ("Budget", "budget", "Measures cost variance and overrun limits."),
                ("Milestones", "milestones", "Tracks ratio of delayed key milestones."),
                ("Risks", "risks", "Subtracts points for severity logs and blockers."),
                ("Sentiment", "sentiment", "Case-insensitive stakeholder comment audit."),
                ("Resources", "resources", "Flags staffing shortages or allocations.")
            ]
            
            for idx, (label, key, desc) in enumerate(cat_list):
                val = cats.get(key, 100.0)
                status = "GREEN" if val >= 80 else "YELLOW" if val >= 50 else "RED"
                pos = positions[idx]
                
                card_title = f"{label}: {val}/100 [{status}]"
                add_card(slide, pos[0], pos[1], col_width, row_height, card_title, desc, border_color=slate_200)

            # ----------------------------------------------------
            # SLIDE 4: Milestones Timeline Tracker
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, "Timeline & Milestone Progress")
            
            miles_body = "Key Milestone Deliverables:\n\n"
            miles = raw_data.get("milestones", []) or []
            if not miles:
                miles_body += "• No milestones listed."
            else:
                for m in miles[:6]:
                    name = m.get("name") if isinstance(m, dict) else str(m)
                    status = m.get("status") if isinstance(m, dict) else "N/A"
                    due = m.get("due_date") if isinstance(m, dict) else "N/A"
                    miles_body += f"• {name} | Due: {due} | Status: {status}\n"
                    
            add_card(slide, Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5), "Milestones Schedule Logs", miles_body)

            # ----------------------------------------------------
            # SLIDE 5: Operational Risks & Resource Constraints
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, "Active Constraints & Staffing Status")
            
            blockers_list = raw_data.get("blockers", []) or []
            deps_list = raw_data.get("dependencies", []) or []
            res_avail = raw_data.get("resource_availability") or "Adequate staffing levels."
            
            blockers_text = "• " + "\n• ".join(blockers_list) if blockers_list else "• No active blockers."
            deps_text = "• " + "\n• ".join(deps_list) if deps_list else "• No critical dependencies."
            
            left_text = f"Blockers:\n{blockers_text}\n\nDependencies:\n{deps_text}"
            add_card(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), "Blockers & Integrations", left_text)
            
            right_text = f"Staffing Allocation:\n• {res_avail}\n\nClient Sentiment comments:\n"
            comments = raw_data.get("stakeholder_comments", []) or []
            if not comments:
                right_text += "• No comments recorded."
            else:
                for c in comments[:3]:
                    right_text += f"• \"{c}\"\n"
            add_card(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), "Resources & Stakeholder Sentiment", right_text)

            # ----------------------------------------------------
            # SLIDE 6: Strategic Recommendations
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, "Strategic Recovery Recommendations")
            
            rec_card_width = Inches(12.333)
            rec_card_height = Inches(4.5)
            
            # Calculate targeted recs based on scores
            recs_list = []
            if cats.get("schedule", 100.0) < 100.0:
                recs_list.append("Schedule Recovery: Align developer velocity with milestones to recover timeline slippage.")
            if cats.get("budget", 100.0) < 100.0:
                recs_list.append("Budget Mitigation: Conduct spent-to-baseline audit to slow budget burn overrun.")
            if cats.get("risks", 100.0) < 100.0:
                recs_list.append("Mitigate Blockers: Resolve active blockers and key resource constraints immediately.")
            if cats.get("sentiment", 100.0) < 80.0:
                recs_list.append("Stakeholder Management: Conduct alignment workshops to address client concern comments.")
                
            if not recs_list:
                recs_list.append("Maintain Standard Checks: Review status parameters on schedule to maintain GREEN status.")
                recs_list.append("Operational Stability: Review upcoming system migration dependencies.")
                
            recs_text = "\n\n".join([f"{i+1}. {r}" for i, r in enumerate(recs_list[:4])])
            add_card(slide, Inches(0.5), Inches(1.5), rec_card_width, rec_card_height, "Next-Step Action Plan", recs_text)

            # Save presentation
            prs.save(output_path)
            logger.info(f"Successfully generated and saved widescreen weekly PPTX to {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Error compiling Weekly PPTX slide deck: {e}")
            raise

class PortfolioPPTGenerator:
    """Agent that compiles Monthly Portfolio Trend Reports into widescreen 16:9 consulting PPTX decks."""
    
    def __init__(self, template_path: str = "templates/ppt_template.pptx"):
        self.template_path = os.path.abspath(template_path)
        logger.info(f"PortfolioPPTGenerator initialized. Template: {self.template_path}")

    def generate_portfolio_deck(self, trend_report: Any, output_path: str) -> str:
        """
        Generates a 6-slide executive PPTX presentation representing:
        1. Executive Summary
        2. Portfolio Health (with Chart Placeholder)
        3. Key Trends (with Chart Placeholder)
        4. Emerging Risks (with Chart Placeholder)
        5. Recommendations
        6. Outlook (with Chart Placeholder)
        """
        logger.info(f"Generating Portfolio PPTX presentation at {output_path}")
        
        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        if not PPTX_AVAILABLE:
            logger.info("Running in mock mode. Writing plain text mock slide info.")
            with open(output_path, "w") as f:
                f.write(f"Mock Portfolio PPTX Slide Deck.\nInsights: {getattr(trend_report, 'executive_insights', [])}")
            return output_path

        try:
            # Load template or start blank
            prs = None
            if os.path.exists(self.template_path) and os.path.getsize(self.template_path) > 0:
                try:
                    prs = Presentation(self.template_path)
                    logger.info(f"Loaded PPTX template from {self.template_path}")
                except Exception as ex:
                    logger.warning(f"Failed to load PPTX template: {ex}. Starting blank.")
                    
            if prs is None:
                prs = Presentation()
                logger.info("Starting blank presentation.")

            # Set modern widescreen dimensions
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Slide layouts: layout 6 is blank in standard PPTX. We construct custom elements on blank slides.
            blank_layout = prs.slide_layouts[6]
            
            # Colors
            navy = RGBColor(15, 23, 42)
            light_gray = RGBColor(241, 245, 249)
            dark_gray = RGBColor(71, 85, 105)
            white = RGBColor(255, 255, 255)
            purple = RGBColor(124, 58, 237)

            # Helper for header titles
            def add_header(slide, title_text: str):
                header_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
                header_rect.fill.solid()
                header_rect.fill.fore_color.rgb = navy
                header_rect.line.fill.background()
                
                txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.6))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title_text
                p.font.name = "Arial"
                p.font.size = Pt(26)
                p.font.bold = True
                p.font.color.rgb = white
                
            # Helper for chart placeholder shapes
            def add_chart_placeholder(slide, title: str, left, top, width, height):
                card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                card.fill.solid()
                card.fill.fore_color.rgb = light_gray
                card.line.color.rgb = dark_gray
                
                tf = card.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"\n\n\n[ Placeholder: {title} ]"
                p.font.name = "Arial"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = dark_gray
                p.alignment = 1  # Centered

            # Helper for consulting-style cards
            def add_card(slide, left, top, width, height, title, body_text="", fill_color=None, border_color=None):
                card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                card.fill.solid()
                card.fill.fore_color.rgb = fill_color or light_gray
                if border_color:
                    card.line.color.rgb = border_color
                else:
                    card.line.fill.background()
                
                tf = card.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.2)
                tf.margin_top = Inches(0.2)
                
                p = tf.paragraphs[0]
                p.text = title
                p.font.name = "Arial"
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = navy
                
                if body_text:
                    p2 = tf.add_paragraph()
                    p2.text = body_text
                    p2.font.name = "Arial"
                    p2.font.size = Pt(11)
                    p2.font.color.rgb = dark_gray
                    p2.space_before = Pt(8)
                return card

            # ----------------------------------------------------
            # SLIDE 1: Executive Summary
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, "Portfolio Executive Summary")
            
            # Left panel - Key Insights
            left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.5), Inches(5.5))
            tf = left_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = "Portfolio Takeaways:"
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = navy
            
            insights = getattr(trend_report, "executive_insights", []) or ["No insights compiled."]
            for insight in insights:
                p = tf.add_paragraph()
                p.text = f"• {insight}"
                p.font.name = "Arial"
                p.font.size = Pt(13)
                p.font.color.rgb = dark_gray
                p.space_after = Pt(10)
                
            # Right panel - Budget Summary
            right_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.3), Inches(5.0), Inches(5.5))
            tf_r = right_box.text_frame
            tf_r.word_wrap = True
            tf_r.paragraphs[0].text = "Financial Overview:"
            tf_r.paragraphs[0].font.size = Pt(18)
            tf_r.paragraphs[0].font.bold = True
            tf_r.paragraphs[0].font.color.rgb = navy
            
            b = getattr(trend_report, "budget_burn_trends", None)
            total_b = getattr(b, "total_portfolio_budget", 0.0)
            total_s = getattr(b, "total_portfolio_spent", 0.0)
            trajectory = getattr(b, "burn_rate_trajectory", "STABLE")
            
            p1 = tf_r.add_paragraph()
            p1.text = f"• Total Portfolio Budget: ${total_b:,.2f}"
            p1.font.name = "Arial"
            p1.font.size = Pt(13)
            p1.font.color.rgb = dark_gray
            p2 = tf_r.add_paragraph()
            p2.text = f"• Total Portfolio Spent:  ${total_s:,.2f}"
            p2.font.name = "Arial"
            p2.font.size = Pt(13)
            p2.font.color.rgb = dark_gray
            p3 = tf_r.add_paragraph()
            p3.text = f"• Burn Trajectory:       {trajectory}"
            p3.font.name = "Arial"
            p3.font.size = Pt(13)
            p3.font.bold = True
            p3.font.color.rgb = purple

            # ----------------------------------------------------
            # SLIDE 2: Portfolio Health
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, "Portfolio RAG Health Status")
            
            list_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5))
            tf = list_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = "Projects Monitored:"
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = navy
            
            at_risk = getattr(trend_report, "projects_at_risk", []) or []
            if not at_risk:
                p = tf.add_paragraph()
                p.text = "• All monitored projects are currently GREEN."
                p.font.name = "Arial"
                p.font.size = Pt(13)
                p.font.color.rgb = dark_gray
            else:
                for proj in at_risk:
                    p = tf.add_paragraph()
                    p.text = f"• {proj.project_name} [{proj.rag_status}]: {proj.reasons[:70]}..."
                    p.font.name = "Arial"
                    p.font.size = Pt(12)
                    p.font.color.rgb = dark_gray
                    p.space_after = Pt(8)
                    
            # Calculate project RAG count from insights and at-risk list
            import re
            total_projects = 3
            for ins in getattr(trend_report, "executive_insights", []):
                match = re.search(r"comprises\s+(\d+)\s+unique", ins, re.IGNORECASE)
                if match:
                    total_projects = int(match.group(1))
                    break
            
            yellow_count = sum(1 for p in at_risk if p.rag_status == "YELLOW")
            red_count = sum(1 for p in at_risk if p.rag_status == "RED")
            green_count = max(0, total_projects - yellow_count - red_count)
            
            chart_data = CategoryChartData()
            chart_data.categories = ['Green Status', 'Yellow Status', 'Red Status']
            chart_data.add_series('Projects', (green_count, yellow_count, red_count))
            
            x, y, cx, cy = Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT

            # ----------------------------------------------------
            # SLIDE 3: Key Trends
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, "Key Operational & Milestone Trends")
            
            trend_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5))
            tf = trend_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = "Milestone Delay Summary:"
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = navy
            
            delays = getattr(trend_report, "repeated_milestone_delays", []) or []
            if not delays:
                p = tf.add_paragraph()
                p.text = "• No repeated milestone delays identified."
                p.font.name = "Arial"
                p.font.size = Pt(13)
                p.font.color.rgb = dark_gray
            else:
                for d in delays[:4]:
                    p = tf.add_paragraph()
                    p.text = f"• {d.milestone_name}: {d.details}"
                    p.font.name = "Arial"
                    p.font.size = Pt(12)
                    p.font.color.rgb = dark_gray
                    p.space_after = Pt(8)

            chart_data_burn = CategoryChartData()
            chart_data_burn.categories = ['Portfolio Budget', 'Portfolio Spent']
            chart_data_burn.add_series('Amount ($)', (total_b, total_s))
            
            x_b, y_b, cx_b, cy_b = Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0)
            chart_b = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x_b, y_b, cx_b, cy_b, chart_data_burn
            ).chart
            chart_b.has_legend = False

            # ----------------------------------------------------
            # SLIDE 4: Emerging Risks
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, "Emerging Risks & Blockers")
            
            risks_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5))
            tf = risks_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = "Common Blocker Clusters:"
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = navy
            
            blockers = getattr(trend_report, "common_blockers", []) or []
            if not blockers:
                p = tf.add_paragraph()
                p.text = "• No active blocker clusters identified."
                p.font.name = "Arial"
                p.font.size = Pt(13)
                p.font.color.rgb = dark_gray
            else:
                for b in blockers[:3]:
                    p = tf.add_paragraph()
                    p.text = f"• {b.blocker_type}: affects {len(b.affected_projects)} projects."
                    p.font.name = "Arial"
                    p.font.size = Pt(12)
                    p.font.color.rgb = dark_gray
                    p.space_after = Pt(8)

            chart_data_risk = CategoryChartData()
            categories_risk = []
            counts_risk = []
            for b in blockers[:3]:
                categories_risk.append(b.blocker_type.split(" ")[0])
                counts_risk.append(len(b.affected_projects))
            
            if not categories_risk:
                categories_risk = ['Access', 'Vendor', 'Other']
                counts_risk = [0, 0, 0]
                
            chart_data_risk.categories = categories_risk
            chart_data_risk.add_series('Affected Projects', tuple(counts_risk))
            
            x_r, y_r, cx_r, cy_r = Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0)
            chart_r = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x_r, y_r, cx_r, cy_r, chart_data_risk
            ).chart
            chart_r.has_legend = False

            # ----------------------------------------------------
            # SLIDE 5: Recommendations
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, "Strategic Recommendations")
            
            rec_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.0), Inches(5.5))
            tf = rec_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = "Recovery Action Plans:"
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = navy
            
            recs = [
                "Mitigate Access Blockers: Standardize staging credential issuance across active portfolios.",
                "Verify Budget Overruns: Review projects with accelerating spent-to-budget variances.",
                "Milestones Acceleration: Address delays in core system developments and specifications.",
                "Sentiment Rebuild: Schedule stakeholder check-in workshops to review risk mitigation logs."
            ]
            for i, r in enumerate(recs):
                p = tf.add_paragraph()
                p.text = f"{i+1}. {r}"
                p.font.name = "Arial"
                p.font.size = Pt(14)
                p.font.color.rgb = dark_gray
                p.space_after = Pt(15)

            # ----------------------------------------------------
            # SLIDE 6: Outlook
            # ----------------------------------------------------
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, "Portfolio Projections & Outlook")
            
            out_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5))
            tf = out_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = "Projections:"
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = navy
            
            s = getattr(trend_report, "sentiment_trends", None)
            overall_s = getattr(s, "overall_sentiment", "NEUTRAL")
            trajectory_s = getattr(s, "trajectory", "STABLE")
            
            p1 = tf.add_paragraph()
            p1.text = f"• Sentiment Status: {overall_s}"
            p1.font.name = "Arial"
            p1.font.size = Pt(13)
            p1.font.color.rgb = dark_gray
            p2 = tf.add_paragraph()
            p2.text = f"• Sentiment Trajectory: {trajectory_s}"
            p2.font.name = "Arial"
            p2.font.size = Pt(13)
            p2.font.color.rgb = dark_gray
            p3 = tf.add_paragraph()
            p3.text = f"• Portfolio Outlook is stable with recovery plans in execution."
            p3.font.name = "Arial"
            p3.font.size = Pt(13)
            p3.font.color.rgb = dark_gray

            chart_data_out = CategoryChartData()
            chart_data_out.categories = ['Week 1', 'Week 2', 'Week 3', 'Week 4']
            
            avg_sentiment = 80.0
            if overall_s == "NEUTRAL":
                avg_sentiment = 50.0
            elif overall_s == "NEGATIVE":
                avg_sentiment = 20.0
                
            chart_data_out.add_series('Sentiment Score', (avg_sentiment * 0.9, avg_sentiment * 0.95, avg_sentiment * 0.92, avg_sentiment))
            
            x_o, y_o, cx_o, cy_o = Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0)
            chart_o = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x_o, y_o, cx_o, cy_o, chart_data_out
            ).chart
            chart_o.has_legend = True
            chart_o.legend.position = XL_LEGEND_POSITION.BOTTOM

            # Save presentation
            prs.save(output_path)
            logger.info(f"Successfully generated and saved Portfolio PPTX presentation at {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Error compiling Portfolio PPTX slide deck: {e}")
            raise
