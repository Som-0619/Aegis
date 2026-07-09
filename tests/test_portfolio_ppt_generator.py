"""
Unit tests for the PortfolioPPTGenerator.
"""

import pytest
import os
from agents.ppt_generator import PortfolioPPTGenerator, PPTX_AVAILABLE
from agents.portfolio_analyzer import PortfolioTrendReport, RepeatedMilestoneDelay, BudgetBurnTrends, CommonBlocker, ProjectRiskProfile, SentimentTrends

def test_portfolio_ppt_generator_instantiation():
    """Verifies that PortfolioPPTGenerator can be instantiated."""
    generator = PortfolioPPTGenerator()
    assert generator is not None

def test_portfolio_deck_generation(tmp_path):
    """
    Tests compiling a 6-slide executive PPTX deck from a mock trend report.
    Verifies slide count and content assertions if pptx is available.
    """
    generator = PortfolioPPTGenerator()
    
    # 1. Mock trend report
    mock_report = PortfolioTrendReport(
        repeated_milestone_delays=[
            RepeatedMilestoneDelay(milestone_name="QA Testing", affected_projects=["Project Alpha"], details="Delayed due to environment issues")
        ],
        budget_burn_trends=BudgetBurnTrends(
            total_portfolio_budget=500000.0,
            total_portfolio_spent=350000.0,
            burn_rate_trajectory="STABLE",
            insights="Budget is within bounds."
        ),
        common_blockers=[
            CommonBlocker(blocker_type="Credentials Blocker", affected_projects=["Project Alpha"], details="Access key delivery delay")
        ],
        projects_at_risk=[
            ProjectRiskProfile(project_name="Project Alpha", rag_status="YELLOW", reasons="Milestone delays")
        ],
        improving_projects=[],
        sentiment_trends=SentimentTrends(
            overall_sentiment="NEUTRAL",
            trajectory="STABLE",
            details="Stakeholders are neutral"
        ),
        executive_insights=[
            "Portfolio health is stable.",
            "Review credentials issue."
        ]
    )
    
    output_pptx = tmp_path / "portfolio_trends.pptx"
    
    # 2. Compile deck
    res_path = generator.generate_portfolio_deck(mock_report, str(output_pptx))
    
    assert os.path.exists(res_path)
    
    # 3. Assert slide count and structures
    if PPTX_AVAILABLE:
        from pptx import Presentation  # type: ignore
        prs = Presentation(res_path)
        
        # Verify it has exactly 6 slides
        assert len(prs.slides) == 6
        
        # Verify title headers on each slide
        expected_titles = [
            "Portfolio Executive Summary",
            "Portfolio RAG Health Status",
            "Key Operational & Milestone Trends",
            "Emerging Risks & Blockers",
            "Strategic Recommendations",
            "Portfolio Projections & Outlook"
        ]
        
        for idx, slide in enumerate(prs.slides):
            # Check shapes to find the title text
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        texts.append(paragraph.text)
            
            # Title should be the expected slide title
            assert any(expected_titles[idx] in t for t in texts), f"Slide {idx+1} missing title: {expected_titles[idx]}"
            
            # Check for charts on slides 2, 3, 4, 6
            if idx in (1, 2, 3, 5):
                assert any(shape.has_chart for shape in slide.shapes), f"Slide {idx+1} missing chart shape."
